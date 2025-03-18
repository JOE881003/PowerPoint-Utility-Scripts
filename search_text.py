from pptx import Presentation
import os
           
def ppt_search_text(file_path, text):
    n = 0
    prs = Presentation(os.path.abspath(file_path))
    for slide in prs.slides:
        n += 1
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.lower()
                if str(text) in shape.text:
                    print(shape.text)
                    print('=================================')


path = "t.pptx"
text = "人工智慧"
ppt_search_text(path, text)